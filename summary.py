import os
import json
from loguru import logger
from core.meta_oracle import get_session
from dao.file_orm import FileSummary
from dao.model_repo import KbotMdModelsRepository
from dao.file_repo import FileSummaryRepository
from oci_client import OCIClient, OCILLMConfig
from core.settings import get_app_config

import fitz  # PyMuPDF
from docx import Document
import openpyxl
from pptx import Presentation



class FileSummaryService:

    @property
    def db_session(self):
        return get_session()

    async def summary(self, file_ids: list[str]):
        """文件摘要服务"""
        async with self.db_session as session:
            file_repo = FileSummaryRepository(db_session=session)
            files = await file_repo.get_by_ids(file_ids)
            if not files:
                raise Exception(f"文件 {file_ids} 不存在")
            
            # 获取模型客户端
            client = await self._get_model()
            # 启动模型
            await client.startup()
            # 执行摘要提取
            for file in files:
                # 获取文件内容
                try:
                    file_content = self._get_file_content(file.file_path)
                except Exception as e:
                    logger.error(f"解析文件 {file.file_path} 失败: {e}")
                    continue
                
                # 获取摘要提取提示
                prompt = self._get_prompt(file_content)
                # 执行摘要提取
                response = await client.chat(prompt, stream=True)
                
                # 收集完整的流式响应内容
                full_content = ""
                try:
                    # OCI SDK 特有的事件流处理
                    for event in response.data.events():  # type: ignore
                        event_data = json.loads(event.data)
                        
                        # 根据参考的LLM代码，处理不同格式的响应
                        text = None
                        
                        # 1. OCI Cohere 格式: {"apiFormat": "COHERE", "text": "你好", "pad": "..."}
                        if 'apiFormat' in event_data and event_data.get('apiFormat') == 'COHERE':
                            text = event_data.get('text', '')
                        
                        # 2. OCI Generic/Grok 格式: {"index": 0, "message": {"role": "ASSISTANT", "content": [{"type": "TEXT", "text": "你好"}]}, "pad": "..."}
                        elif 'message' in event_data and isinstance(event_data.get('message'), dict):
                            message = event_data['message']
                            content = message.get('content', [])
                            if content and isinstance(content, list) and len(content) > 0:
                                # 提取 content[0].text
                                first_content = content[0]
                                if isinstance(first_content, dict) and first_content.get('type') == 'TEXT':
                                    text = first_content.get('text', '')
                        
                        # 3. 直接包含文本内容
                        elif 'text' in event_data:
                            text = event_data.get('text', '')
                        
                        # 4. 如果包含choices字段（OpenAI兼容格式）
                        elif 'choices' in event_data and isinstance(event_data['choices'], list) and len(event_data['choices']) > 0:
                            choice = event_data['choices'][0]
                            if 'delta' in choice and 'content' in choice['delta']:
                                text = choice['delta']['content']
                        
                        if text is not None:
                            full_content += str(text)
                        
                except Exception as e:
                    logger.exception(f"OCI流式响应错误: {e}")
                    raise
                
                # 将完整内容作为结果
                result = full_content
                
                # 处理结果，提取摘要信息
                try:
                    if result:
                        try:
                            # 尝试解析JSON格式的响应
                            summary_data = json.loads(result)
                            # 验证是否包含所有必要的语言版本
                            required_langs = ["cn", "en", "ja", "ko"]
                            if all(lang in summary_data for lang in required_langs):
                                # 更新文件摘要
                                kwargs = {
                                    "summary_cn": summary_data.get("cn", ""),
                                    "summary_en": summary_data.get("en", ""),
                                    "summary_ja": summary_data.get("ja", ""),
                                    "summary_kr": summary_data.get("ko", ""),
                                    "status": 2
                                }
                                await file_repo.update(file.file_id, **kwargs)
                                logger.info(f"文件 {file.file_name} 摘要生成成功")
                            else:
                                logger.warning(f"文件 {file.file_name} 摘要格式不完整: {result}")
                        except json.JSONDecodeError:
                            # 如果不是JSON格式，解析多语言Markdown内容
                            summaries = self._parse_multilingual_summary(result)
                            kwargs = {
                                "summary_cn": summaries.get("cn", ""),
                                "summary_en": summaries.get("en", ""),
                                "summary_ja": summaries.get("ja", ""),
                                "summary_kr": summaries.get("ko", ""),
                                "status": 2
                            }
                            await file_repo.update(file.file_id, **kwargs)
                            logger.info(f"文件 {file.file_name} 摘要生成成功（多语言Markdown格式）")
                except Exception as e:
                    logger.error(f"文件 {file.file_name} 摘要生成失败: {e}")
                    kwargs = {"status": 3}
                    await file_repo.update(file.file_id, **kwargs)

    def _get_file_content(self, file_path):
        ext = file_path.split('.')[-1].lower()
        
        # 1. 处理 PDF
        if ext == 'pdf':
            text = ""
            with fitz.open(file_path) as doc:
                for page in doc:
                    text += page.get_text()
            return text.strip()

        # 2. 处理 Word (.docx)
        elif ext == 'docx':
            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs]).strip()

        # 3. 处理 Excel (.xlsx)
        elif ext == 'xlsx':
            wb = openpyxl.load_model(file_path, data_only=True)
            content = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    content.append(" ".join([str(cell) for cell in row if cell is not None]))
            return "\n".join(content).strip()

        # 4. 处理 PPT (.pptx)
        elif ext == 'pptx':
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs).strip()
        
        # 5. 默认处理文本 (txt, md)
        else:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read().strip()

    async def _get_model(self) -> OCIClient:
        """获取模型实例"""
        model_name = get_app_config().llm_model
        # 根据模型名称从数据库获取模型配置信息
        async with self.db_session as session:
            model_repo = KbotMdModelsRepository(db_session=session)
            models = await model_repo.get(model_name)
            if not models:
                raise Exception(f"模型 {model_name} 不存在")
            model = models[0]

        # 解析模型配置信息
        endpoint = model.api_endpoint
        model_param = model.model_params
        if not endpoint or not model_param:
            raise Exception(f"模型 {model_name} 配置不完整")
        config_file = model_param.get("config_file")
        compartment_id = model_param.get("compartment_id")
        if not config_file or not compartment_id:
            raise Exception(f"模型 {model_name} 配置文件不存在")
        
        # 创建 Pydantic 模型实例时传入所有必需参数
        model_config = OCILLMConfig(
            model_name=model.model_name,
            api_endpoint=endpoint,
            config_file=config_file,
            compartment_id=compartment_id,
            max_tokens=int(model_param.get("max_tokens", 8192)),
            temperature=float(model_param.get("temperature")) if model_param.get("temperature") else None
        )

        # 创建模型客户端
        client = OCIClient(model_config)
        return client
    
    async def _get_summary(self, client: OCIClient, prompt: str):
        """执行摘要提取"""

        return await client.chat(prompt)
    

    def _parse_multilingual_summary(self, content: str) -> dict[str, str]:
        """解析多语言Markdown摘要内容"""
        summaries = {"cn": "", "en": "", "ja": "", "ko": ""}
        
        # 定义语言标识符
        language_headers = {
            "cn": ["# 目录", "# 摘要"],
            "en": ["# Table of Contents", "# Summary"],
            "ja": ["# 目次", "# 要約"],
            "ko": ["# 목차", "# 요약"]
        }
        
        current_lang = None
        lines = content.split('\n')
        current_content = []
        
        for line in lines:
            # 检查是否为语言标题
            for lang, headers in language_headers.items():
                if any(line.strip().startswith(header) for header in headers):
                    # 保存上一个语言的内容
                    if current_lang and current_content:
                        summaries[current_lang] = '\n'.join(current_content).strip()
                    # 开始新的语言
                    current_lang = lang
                    current_content = [line]
                    break
            else:
                # 如果不是标题行，添加到当前语言内容
                if current_lang:
                    current_content.append(line)
        
        # 保存最后一个语言的内容
        if current_lang and current_content:
            summaries[current_lang] = '\n'.join(current_content).strip()
        
        return summaries

    def _get_prompt(self, file_content: str) -> str:
        """获取摘要提取提示"""
        prompt_path = os.path.join(os.path.dirname(__file__), "prompt.txt")
        try:
            with open(prompt_path, 'r', encoding='utf-8') as f:
                prompt_content = f.read().strip()
                return prompt_content + file_content
        except Exception as e:
            logger.error(f"读取 prompt.txt 失败: {e}")
            prompt_content = f"""
Please process the document I provide in strict accordance with the following requirements, and return only the final dictionary without any extra explanation or comment:
Extract the core table of contents from the document (keep it concise, first-level sections only).
Generate a concise and accurate summary of the document.
Combine the table of contents and summary into Markdown format, structured as:
Table of Contents
...
Summary
...
Produce four versions: Chinese (cn), English (en), Japanese (ja), Korean (ko).
Return the result only as a dictionary in the following format:
{
"cn": "Chinese Markdown content",
"en": "English Markdown content",
"ja": "Japanese Markdown content",
"ko": "Korean Markdown content"
}
Rules:
Translate naturally and accurately, following local language habits.
Do not add subjective opinions or extra information.
Strictly follow the format; do not change keys or structure.
Now process the document:

{file_content}"""
            return prompt_content